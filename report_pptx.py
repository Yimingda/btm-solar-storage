"""
report_pptx.py  —  Executive PPTX Report Generator
BTM PV+BESS Financial Modelling System

Design language  (Huawei Digital Power style)
──────────────────────────────────────────────
• White background slides with near-black (#0D0D0D) header bars
• Huawei Red (#CF0A2C) — structural accents: header strips, chart bars,
  section markers, left-border strips
• Huawei Gold (#FFC000) — data / KPI values (aligned with Huawei template)
• Font: Calibri (closest widely-available substitute for Huawei Sans)
• McKinsey "pyramid" headline structure: slide title = the conclusion
• 2-3 sentence narrative paragraph per content slide
• Footer: EPC company name  ·  Solution Info URL  ·  page number

Slide sequence  (11 slides)
────────────────────────────
  1.  Cover
  2.  Investment Thesis
  3.  System Overview
  4.  Financial Returns
  5.  Energy Analysis
  6.  Tariff Opportunity
  7.  Implementation Roadmap
  8.  Huawei Digital Power — Partner Credentials  ← NEW
  9.  South Africa Reference Projects              ← NEW
  10. SA Market Context (2026 industry data)
  11. Assumptions & Disclaimer

Solution Info: https://info.support.huawei.com/Energy/info/en_US/all/index
"""
from __future__ import annotations

import io
import os
import sys
from datetime import date

# ── Asset manifest (PV / BESS scenario-tagged media catalogue) ────────────────
_SELF_DIR = os.path.dirname(os.path.abspath(__file__))
if _SELF_DIR not in sys.path:
    sys.path.insert(0, _SELF_DIR)
try:
    import assets_manifest as _am
except ImportError:
    _am = None   # type: ignore[assignment]

# ── Palette (Huawei Digital Power) ───────────────────────────────────────────
_BLK  = "0D0D0D"   # near-black   — header bars, cover background
_HRD  = "CF0A2C"   # Huawei red   — structural accents, chart bars, left strips
_GLD  = "FFC000"   # Huawei gold  — KPI values, capacity specs (matches template)
_DRD  = "9B0821"   # dark red     — heavy emphasis
_LRD  = "FFF5F6"   # pale red     — card/block backgrounds
_DGY  = "1A1A1A"   # dark grey    — slide titles, body text
_MGY  = "6B7280"   # mid grey     — labels, captions, subtitles
_LGRY = "F5F5F5"   # light grey   — alternating row tint
_AMBR = "F59E0B"   # amber        — timeline operational phase
_NEG  = "DC2626"   # bright red   — negative cash-flow bars
_WHT  = "FFFFFF"   # white
_SEP  = "E5E7EB"   # light separator
_DNAV = "111827"   # dark navy    — credential / project card backgrounds

_FONT = "Calibri"   # Huawei Sans substitute

_SOL_URL  = "https://info.support.huawei.com/Energy/info/en_US/all/index"
_SOL_LBL  = "Solution Info"

# ── Auto-scale unit helper ────────────────────────────────────────────────────
_MW_UP = {"kWp": "MWp", "kWh": "MWh", "kW": "MW", "kWh/yr": "MWh/yr"}

def _fmw(v: float, unit: str = "kWp") -> str:
    """Return formatted string: auto-upgrades to MW/MWh/MWp when v ≥ 1 000."""
    if v >= 1000:
        return f"{v / 1000:,.1f} {_MW_UP.get(unit, unit)}"
    return f"{v:,.0f} {unit}"

# ── python-pptx micro-helpers ─────────────────────────────────────────────────

def _rgb(h: str):
    from pptx.dml.color import RGBColor
    h = h.lstrip("#")
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

def _in(v: float):
    from pptx.util import Inches; return Inches(v)

def _pt(v: float):
    from pptx.util import Pt; return Pt(v)

def _rect(slide, x, y, w, h, fill: str, line: str | None = None):
    shp = slide.shapes.add_shape(1, _in(x), _in(y), _in(w), _in(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
    if line:
        shp.line.color.rgb = _rgb(line); shp.line.width = _pt(0.75)
    else:
        shp.line.fill.background()
    return shp

def _tb(slide, x, y, w, h, text: str, size: float,
        bold=False, color=_DGY, align="left", wrap=True):
    from pptx.enum.text import PP_ALIGN
    _am = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
           "right": PP_ALIGN.RIGHT}
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = _pt(0)
    para = tf.paragraphs[0]; para.alignment = _am.get(align, PP_ALIGN.LEFT)
    run = para.add_run(); run.text = text
    run.font.name = _FONT
    run.font.size = _pt(size); run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return tb

# ── Shared layout elements ────────────────────────────────────────────────────

def _header_bar(slide, title: str, subtitle: str = ""):
    """White slide — thin Huawei-black top bar + red left border marker."""
    _rect(slide, 0, 0, 13.33, 0.52, _BLK)          # black top bar
    _rect(slide, 0, 0, 0.09,  0.52, _HRD)           # red left strip on bar
    _rect(slide, 0, 0.52, 13.33, 6.98, _WHT)        # white body
    _rect(slide, 0, 0.52, 0.09, 6.98, _LRD)         # faint red left border

    _tb(slide, 0.28, 0.08, 12.8, 0.40, title,
        19, bold=True, color=_WHT)
    if subtitle:
        _tb(slide, 0.28, 0.55, 12.8, 0.35, subtitle, 9.5, color=_MGY)

def _footer(slide, company: str, page: int, total: int = 11):
    _rect(slide, 0, 7.22, 13.33, 0.02, _SEP)
    _tb(slide, 0.28, 7.26, 4.5, 0.22, company, 7.5, color=_MGY)
    _tb(slide, 4.8,  7.26, 5.8, 0.22,
        f"Solution Info: {_SOL_URL}", 7.5, color=_HRD)
    _tb(slide, 11.8, 7.26, 1.5, 0.22, f"{page} / {total}",
        7.5, color=_MGY, align="right")

def _kpi_block(slide, x, y, label: str, value: str,
               sub: str = "", w: float = 2.8):
    """KPI card: Huawei gold value, red left strip (aligned with Huawei template)."""
    _rect(slide, x, y, w, 1.38, _LRD)
    _rect(slide, x, y, 0.07, 1.38, _HRD)            # red left accent
    _tb(slide, x+0.16, y+0.10, w-0.22, 0.28,
        label, 8, bold=True, color=_MGY)
    _tb(slide, x+0.16, y+0.38, w-0.22, 0.60,
        value, 25, bold=True, color=_GLD)             # gold value
    if sub:
        _tb(slide, x+0.16, y+1.06, w-0.22, 0.26, sub, 7.5, color=_MGY)

def _narrative(slide, x, y, w, h, text: str):
    _tb(slide, x, y, w, h, text, 10, color=_MGY, wrap=True)

def _section_hdr(slide, x, y, w, h, text: str):
    """Black section header bar with red left strip."""
    _rect(slide, x, y, w, h, _BLK)
    _rect(slide, x, y, 0.07, h, _HRD)
    _tb(slide, x+0.16, y+(h-0.28)/2, w-0.22, 0.28,
        text, 9.5, bold=True, color=_WHT)

# ── Chart helpers (matplotlib → PNG, Huawei red palette) ─────────────────────

def _style_ax(ax):
    ax.set_facecolor("white")
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280", labelsize=8)
    ax.grid(axis="y", linestyle="--", alpha=0.3, zorder=0, color="#E5E7EB")

def _png_cashflow(fin_df) -> bytes | None:
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        years = fin_df["Year"].tolist()
        ncf   = fin_df["Net Cash Flow NCF (ZAR)"].values / 1e6
        cum   = fin_df["Cumulative CF (ZAR)"].values / 1e6
        fig, ax1 = plt.subplots(figsize=(8.8, 3.2), facecolor="white")
        clrs = [f"#{_HRD}" if v >= 0 else f"#{_NEG}" for v in ncf]
        ax1.bar(years, ncf, color=clrs, width=0.65, alpha=0.88, zorder=3)
        ax1.axhline(0, color="#E5E7EB", linewidth=0.8)
        ax1.set_xlabel("Year", fontsize=9, color="#6B7280", fontname=_FONT)
        ax1.set_ylabel("Annual NCF  (R million)", fontsize=9,
                       color="#6B7280", fontname=_FONT)
        _style_ax(ax1)
        ax2 = ax1.twinx()
        ax2.plot(years, cum, color=f"#{_BLK}", linewidth=2.2,
                 marker="o", markersize=3.5, zorder=5)
        ax2.set_ylabel("Cumulative  (R million)", fontsize=9,
                       color=f"#{_BLK}", fontname=_FONT)
        ax2.tick_params(colors=f"#{_BLK}", labelsize=8)
        ax2.spines[["top","right"]].set_color("#E5E7EB")
        ax2.spines[["left","bottom"]].set_visible(False)
        for i, v in enumerate(cum):
            if v >= 0 and (i == 0 or cum[i-1] < 0):
                ax2.axvline(years[i], color=f"#{_HRD}", linewidth=1.4,
                            linestyle="--", alpha=0.75)
                ax2.text(years[i]+0.18, max(cum)*0.05,
                         f"Break-even\nYr {years[i]}",
                         fontsize=7.5, color=f"#{_HRD}", va="bottom",
                         fontname=_FONT)
                break
        plt.tight_layout(pad=0.4)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
        plt.close(fig); buf.seek(0); return buf.read()
    except Exception:
        return None

def _png_monthly(pvgis_data: dict) -> bytes | None:
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        months  = ["Jan","Feb","Mar","Apr","May","Jun",
                   "Jul","Aug","Sep","Oct","Nov","Dec"]
        monthly = pvgis_data.get("monthly_kwh", [0]*12)
        fig, ax = plt.subplots(figsize=(8.8, 3.0), facecolor="white")
        bar_clr = [f"#{_DRD}" if i in (5,6,7) else f"#{_HRD}"
                   for i in range(12)]
        bars = ax.bar(months, monthly, color=bar_clr, alpha=0.88,
                      width=0.65, zorder=3)
        ax.set_ylabel("kWh / month", fontsize=9, color="#6B7280",
                      fontname=_FONT)
        _style_ax(ax)
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x()+bar.get_width()/2, h+max(monthly)*0.01,
                    f"{h/1000:.1f}k", ha="center", va="bottom",
                    fontsize=7, color="#6B7280", fontname=_FONT)
        ax.axvspan(4.45, 7.55, color=f"#{_LRD}", alpha=0.9, zorder=0)
        ax.text(6.0, max(monthly)*0.92, "Peak-tariff\nseason",
                ha="center", fontsize=7.5, color=f"#{_DRD}",
                style="italic", fontname=_FONT)
        plt.tight_layout(pad=0.4)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
        plt.close(fig); buf.seek(0); return buf.read()
    except Exception:
        return None

def _png_megaflex_tariff() -> bytes | None:
    """Bar chart — SA Eskom Megaflex blended tariff 2016–2030 (actual + projection)."""
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches

        # ── NERSA-approved Eskom tariff escalations → representative blended rate ─
        # Starting base: ~R 0.90/kWh blended effective (VAT incl.) in FY2015
        _base = 0.90
        _increases = {
            2016: 9.40, 2017: 2.20, 2018: 5.20, 2019: 9.40, 2020: 8.10,
            2021: 15.06, 2022: 9.61, 2023: 18.65, 2024: 12.74, 2025: 11.90,
            2026: 8.76,
        }
        hist_yrs, hist_rates, hist_pcts = [], [], []
        r = _base
        for yr, pct in sorted(_increases.items()):
            r = r * (1 + pct / 100)
            hist_yrs.append(yr)
            hist_rates.append(round(r, 3))
            hist_pcts.append(pct)

        # Projection: FY2027 NERSA approved +8.83%, then ~8.5% p.a.
        proj_yrs, proj_rates = [], []
        _fwd = {2027: 8.83, 2028: 8.50, 2029: 8.50, 2030: 8.50}
        rp = hist_rates[-1]
        for yr, pct in sorted(_fwd.items()):
            rp = rp * (1 + pct / 100)
            proj_yrs.append(yr)
            proj_rates.append(round(rp, 3))

        all_yrs  = hist_yrs + proj_yrs
        all_rates = hist_rates + proj_rates

        fig, ax = plt.subplots(figsize=(11.0, 3.6), facecolor="white")

        # Actual bars
        b1 = ax.bar(hist_yrs, hist_rates, color=f"#{_HRD}", alpha=0.88,
                    width=0.72, zorder=3, label="Actual (NERSA-approved)")
        # Projected bars (hatched, lighter)
        b2 = ax.bar(proj_yrs, proj_rates, color=f"#{_MGY}", alpha=0.40,
                    width=0.72, zorder=3, hatch="///", label="Projected")

        # Trend line
        ax.plot(all_yrs, all_rates, color=f"#{_BLK}", linewidth=2.0,
                marker="o", markersize=4.5, zorder=5, alpha=0.80)

        # Divider: actual vs projected
        ax.axvline(2026.5, color=f"#{_MGY}", linewidth=1.2,
                   linestyle="--", alpha=0.60)
        ax.text(2026.6, max(all_rates) * 0.72, "Projected →",
                fontsize=7.5, color=f"#{_MGY}", va="bottom",
                style="italic", fontname=_FONT)

        # Annotate large-increase years
        annot_yrs  = {2021: 15.06, 2023: 18.65}
        for yr, pct in annot_yrs.items():
            rate = dict(zip(hist_yrs, hist_rates))[yr]
            ax.text(yr, rate + 0.04, f"+{pct:.0f}%",
                    ha="center", va="bottom", fontsize=8,
                    color=f"#{_HRD}", fontweight="bold", fontname=_FONT)

        # All bars: show rate value above
        for bar, rate in zip(b1, hist_rates):
            ax.text(bar.get_x() + bar.get_width() / 2,
                    rate + 0.01, f"R{rate:.2f}",
                    ha="center", va="bottom", fontsize=6.5,
                    color=f"#{_DGY}", fontname=_FONT)
        for bar, rate in zip(b2, proj_rates):
            ax.text(bar.get_x() + bar.get_width() / 2,
                    rate + 0.01, f"R{rate:.2f}",
                    ha="center", va="bottom", fontsize=6.5,
                    color=f"#{_MGY}", fontname=_FONT)

        ax.set_ylabel("Blended Rate  (R / kWh, VAT incl.)",
                      fontsize=9, color=f"#{_MGY}", fontname=_FONT)
        ax.set_ylim(0, max(all_rates) * 1.22)
        ax.set_xticks(all_yrs)
        ax.set_xticklabels([str(y) for y in all_yrs], rotation=45, ha="right")
        _style_ax(ax)

        # CAGR annotation
        cagr = (hist_rates[-1] / hist_rates[0]) ** (1 / (len(hist_yrs) - 1)) - 1
        ax.text(0.02, 0.96,
                f"10-yr CAGR  +{cagr * 100:.1f}% p.a.  "
                f"(FY2016 → FY2026)",
                transform=ax.transAxes, fontsize=8.5, fontweight="bold",
                color=f"#{_HRD}", va="top", fontname=_FONT,
                bbox=dict(boxstyle="round,pad=0.35", facecolor="white",
                          edgecolor=f"#{_SEP}", alpha=0.92))

        ax.legend(fontsize=8, loc="upper left", bbox_to_anchor=(0.0, 0.82),
                  frameon=True, framealpha=0.9, edgecolor=f"#{_SEP}")

        plt.tight_layout(pad=0.40)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
        plt.close(fig); buf.seek(0); return buf.read()
    except Exception:
        return None

# ── Slide builders ────────────────────────────────────────────────────────────

def _s1_cover(prs, project_name: str, client_name: str,
              consultant_name: str, pv_kwp: float,
              bess_kwh: float, tariff_mode: str,
              has_pv: bool = True, has_bess: bool = True):
    """Slide 1 – Cover  (clean-energy photo background, Huawei accent)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Left panel: 3-zone design ────────────────────────────────────────────
    # Zone A (top  0.00 – 4.10"):  dark navy overlay — text area
    # Zone B (mid  4.10 – 6.38"):  product photo visible (green-framed window)
    # Zone C (bot  6.38 – 7.50"):  dark navy overlay — footer / confidential
    #
    # Image: kv-pc.jpg  3840×1720 px  ratio=2.233
    #   Scaled to panel width 8.90"  →  height=3.987"
    #   Positioned so product optical centre (≈48% from top) aligns with
    #   green-band centre (5.24"):  top_y = 5.24 - 3.987×0.48 = 3.33"
    #   Visible slice: 20%–77% of image (full inverter + BESS body)
    #   Dark BG of photo blends seamlessly with dark navy overlays above/below
    _COVER_BG = "0B1929"          # Huawei deep navy
    _ASSETS   = os.path.join(os.path.dirname(__file__), "assets")
    # Scenario-aware cover hero: BESS-only → ESS container; PV-only → utility
    # PV+BESS → retain the existing dark mixed-product shot (hw_kv_dark.jpg)
    if _am is not None:
        _cover_ph = _am.cover_hero(has_pv=has_pv, has_bess=has_bess) \
                    or os.path.join(_ASSETS, "hw_kv_dark.jpg")
    else:
        _cover_ph = os.path.join(_ASSETS, "hw_kv_dark.jpg")

    if os.path.exists(_cover_ph):
        # Photo — z-order lowest; width fills panel, height keeps native ratio
        _ph_w, _ph_h = 8.90, 3.987          # inches: 8.90 / 2.233 = 3.987
        _ph_top      = 3.33                  # top_y so products centre on green band
        slide.shapes.add_picture(_cover_ph,
                                 _in(0), _in(_ph_top),
                                 _in(_ph_w), _in(_ph_h))
        # Dark overlay — TOP zone (text readable over dark bg)
        _rect(slide, 0, 0,    8.90, 4.12, _COVER_BG)
        # Dark overlay — BOTTOM zone
        _rect(slide, 0, 6.38, 8.90, 1.12, _COVER_BG)
    else:
        _rect(slide, 0, 0, 8.90, 7.5, _COVER_BG)   # fallback solid

    # Thin Huawei-green separators at zone boundaries
    _rect(slide, 0, 4.10, 8.90, 0.05, "00A870")     # top → photo boundary
    _rect(slide, 0, 6.38, 8.90, 0.05, "00A870")     # photo → footer boundary

    # Red vertical accent strip (Huawei signature)
    _rect(slide, 8.88, 0, 0.10, 7.5, _HRD)
    # White right panel
    _rect(slide, 8.98, 0, 4.35, 7.5, _WHT)

    # EPC + OEM — top left (OEM hardcoded per Huawei partner requirement)
    epc_label = f"EPC:  {consultant_name}" if consultant_name else "EPC:  —"
    _tb(slide, 0.42, 0.18, 8.0, 0.32, epc_label, 9,
        bold=True, color="AAAAAA")
    _tb(slide, 0.42, 0.52, 8.0, 0.32,
        "OEM:  Huawei Technologies SA PTY LTD", 9,
        bold=True, color="CCCCCC")

    # Thin Huawei-red rule
    _rect(slide, 0.42, 0.92, 8.0, 0.04, _HRD)

    # Main title — adapts to PV/BESS scenario
    if project_name:
        title = project_name
    elif has_pv and has_bess:
        title = f"{_fmw(pv_kwp, 'kWp')} PV + {_fmw(bess_kwh, 'kWh')} BESS"
    elif has_pv:
        title = f"{_fmw(pv_kwp, 'kWp')} Solar PV System"
    else:
        title = f"{_fmw(bess_kwh, 'kWh')} Battery Energy Storage"
    _tb(slide, 0.42, 1.10, 8.0, 1.80, title, 38, bold=True, color=_WHT)

    # Sub-title — adapts to scenario
    _sub_map = {
        (True,  True):  "BTM Solar PV & Battery Energy Storage — Executive Report",
        (True,  False): "BTM Solar PV System — Executive Report",
        (False, True):  "Battery Energy Storage System (BESS) — Executive Report",
    }
    _tb(slide, 0.42, 2.98, 8.0, 0.48,
        _sub_map.get((has_pv, has_bess),
                     "BTM Energy Solution — Executive Report"),
        12.5, color="BBBBBB")

    # Date
    _tb(slide, 0.42, 3.55, 8.0, 0.38,
        f"Report Date:   {date.today().strftime('%B %Y')}",
        10.5, color="999999")

    # Solution reference — bottom left
    _tb(slide, 0.42, 6.85, 8.0, 0.38,
        f"Solution Info: {_SOL_URL}",
        8.5, color=_HRD)

    # Right panel: SYSTEM SIZE — adapts to scenario
    _tb(slide, 9.22, 1.20, 3.90, 0.38, "SYSTEM SIZE", 9,
        bold=True, color=_MGY, align="center")
    _rect(slide, 9.22, 1.60, 3.78, 0.03, _SEP)

    if has_pv and has_bess:
        _tb(slide, 9.22, 1.70, 3.90, 0.85,
            _fmw(pv_kwp, "kWp"), 40, bold=True, color=_GLD, align="center")
        _tb(slide, 9.22, 2.60, 3.90, 0.40,
            "Solar PV Array", 11, color=_MGY, align="center")
        _rect(slide, 9.22, 3.08, 3.78, 0.03, _SEP)
        _tb(slide, 9.22, 3.18, 3.90, 0.72,
            _fmw(bess_kwh, "kWh"), 32, bold=True, color=_DGY, align="center")
        _tb(slide, 9.22, 3.95, 3.90, 0.38,
            "Battery Storage (BESS)", 10, color=_MGY, align="center")
    elif has_pv:
        _tb(slide, 9.22, 2.00, 3.90, 1.10,
            _fmw(pv_kwp, "kWp"), 44, bold=True, color=_GLD, align="center")
        _tb(slide, 9.22, 3.18, 3.90, 0.50,
            "Solar PV Array", 14, color=_MGY, align="center")
    else:
        _tb(slide, 9.22, 2.00, 3.90, 1.10,
            _fmw(bess_kwh, "kWh"), 40, bold=True, color=_GLD, align="center")
        _tb(slide, 9.22, 3.18, 3.90, 0.50,
            "Battery Storage (BESS)", 14, color=_MGY, align="center")

    if tariff_mode:
        _rect(slide, 9.22, 4.80, 3.78, 0.54, _LRD)
        _rect(slide, 9.22, 4.80, 0.07, 0.54, _HRD)
        _tb(slide, 9.36, 4.88, 3.55, 0.36,
            f"Tariff:  {tariff_mode[:28]}", 8.5, color=_DGY)

    # Confidential — "Only For: {client_name}"
    conf_text = (f"Confidential  ·  Only For: {client_name}"
                 if client_name else
                 "Confidential  ·  For Authorised Recipients Only")
    _tb(slide, 0.42, 7.08, 8.0, 0.30, conf_text, 8, color="999999")


def _s2_thesis(prs, results: dict, params: dict, company: str,
               page: int = 2, total: int = 11):
    """Slide 2 – Investment Thesis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    npv     = results.get("npv", 0) or 0
    irr     = results.get("irr", 0) or 0
    payback = results.get("payback") or 0
    capex   = results.get("total_capex", 0) or 0
    lcoe_d  = results.get("lcoe") or {}   # PV-only LCOE (None if no PV)
    lcos_d  = results.get("lcos") or {}   # BESS-only LCOS (None if no BESS)
    # Lifetime MWh for narrative: PV gen (if PV), else BESS discharge (pure BESS)
    avoided = (lcoe_d.get("total_avoided_mwh") or
               lcos_d.get("total_discharge_mwh") or 0)
    esc     = params.get("tariff_escalation", 8.0)

    def _m(v): return f"R {abs(v)/1e6:.1f}M" if abs(v) >= 1e6 else \
                      f"R {abs(v)/1e3:.0f}k"

    headline = (f"At {irr:.1f}% IRR and {payback:.1f}-year payback, "
                "this investment delivers superior risk-adjusted returns")
    _header_bar(slide, headline,
                "Capital is recovered before Year 6 — the system then generates "
                "free cash flow for 14+ years with no incremental CAPEX.")

    # Build KPI block list — always CAPEX/NPV/IRR/Payback, then LCOE and/or LCOS
    _kpis = [
        ("TOTAL CAPEX",    _m(capex),          "All-in investment"),
        ("20-YEAR NPV",    _m(npv),            "After tax & escalation"),
        ("PROJECT IRR",    f"{irr:.1f}%",      "Unlevered, after-tax"),
        ("SIMPLE PAYBACK", f"{payback:.1f} yr","Grid-savings basis"),
    ]
    if lcoe_d:
        _lcoe_v = lcoe_d.get("lcoe_zar_kwh", 0) or 0
        _kpis.append(("1ST YR LCOE",  f"R {_lcoe_v:.2f}", "PV cost / kWh generated"))
    if lcos_d:
        _lcos_v = lcos_d.get("lcos_zar_kwh", 0) or 0
        _kpis.append(("1ST YR LCOS",  f"R {_lcos_v:.2f}", "BESS cost / kWh discharged"))

    # Distribute KPI blocks evenly across slide width (0.28" to 13.06" = 12.78")
    _n = len(_kpis)
    _blk_w  = 12.78 / _n
    for i, (lbl, val, sub) in enumerate(_kpis):
        _kpi_block(slide, 0.28 + i * _blk_w, 1.68, lbl, val, sub, w=_blk_w - 0.10)

    story = (
        f"The project displaces expensive {params.get('tariff_mode','Eskom Megaflex')} "
        f"peak energy (winter peak R {params.get('w_evening_peak',8.13):.2f}/kWh) "
        "with self-generated solar and stored BESS energy.  "
        f"At {esc:.1f}% p.a. tariff escalation, returns compound strongly over "
        "the 20-year model horizon — making early deployment value-accretive.  "
        "Section 12B (100% Year-1 for assets <1 MW; 50/30/20% over 3 years for "
        "larger systems) materially improves early after-tax cash flows."
    )
    _rect(slide, 0.28, 3.35, 12.78, 0.03, _SEP)
    _narrative(slide, 0.28, 3.50, 12.78, 0.75, story)

    d1 = results.get("dispatch_yr1") or {}
    load   = max(1, d1.get("annual_load_kWh", 1))
    gbuy   = d1.get("annual_grid_buy_kWh", 1)
    self_s = min(95, round((1-gbuy/load)*100, 0))

    for xi, (hdr, body) in enumerate([
        ("Grid Independence",
         f"After commissioning, the site sources ≈{self_s:.0f}% of daytime demand "
         "from owned generation — insulating operations from future utility price "
         "shocks and loadshedding exposure."),
        ("Lifetime Avoided Energy",
         f"Over 20 years the system avoids purchasing {avoided:,.0f} MWh from the "
         "grid, representing a measurable long-run cost reduction and carbon footprint "
         "improvement."),
    ]):
        _rect(slide, 0.28+xi*6.40, 4.42, 6.20, 1.82, _LRD)
        _rect(slide, 0.28+xi*6.40, 4.42, 0.07, 1.82, _HRD)
        _tb(slide, 0.44+xi*6.40, 4.52, 5.95, 0.32,
            hdr, 10, bold=True, color=_DGY)
        _narrative(slide, 0.44+xi*6.40, 4.88, 5.95, 1.25, body)

    _footer(slide, company, page, total=total)


def _s3_system(prs, params: dict, pvgis_data: dict, company: str,
               page: int = 3, total: int = 11,
               has_pv: bool = True, has_bess: bool = True):
    """Slide 3 – System Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    annual_kwh = pvgis_data.get("annual_kwh", 0)
    pv_kwp  = params.get("pv_kwp", 0)
    bess_kh = params.get("bess_kwh", 0)
    spec_yield = annual_kwh / pv_kwp if pv_kwp else 0

    # Headline adapts to scenario
    if has_pv and has_bess:
        hl = (f"{_fmw(pv_kwp,'kWp')} PV + {_fmw(bess_kh,'kWh')} BESS "
              f"— specific yield {spec_yield:,.0f} kWh/kWp/yr")
    elif has_pv:
        hl = f"{_fmw(pv_kwp,'kWp')} Solar PV — specific yield {spec_yield:,.0f} kWh/kWp/yr"
    else:
        c_rate = params.get("c_rate", 0.25)
        hl = (f"{_fmw(bess_kh,'kWh')} BESS — "
              f"{_fmw(bess_kh*c_rate,'kW')} max discharge, "
              f"{params.get('dod',90):.0f}% DoD")
    _header_bar(slide, hl,
                f"Lat {params.get('lat',0):.3f}°  Lon {params.get('lon',0):.3f}°  "
                f"·  Tilt {params.get('tilt',20):.0f}°  ·  Azimuth {params.get('azimuth',180):.0f}°  "
                f"·  Irradiance: EU PVGIS API (Joint Research Centre)")

    # ── Column width ──────────────────────────────────────────────────────────
    # PV+BESS  → two equal 5.95" columns (unchanged dual layout)
    # Single   → spec table 6.60" on the left; product image panel on the right
    _dual = has_pv and has_bess
    _W    = 5.95 if _dual else 6.60
    _x2   = 7.05  # x-position of BESS column in dual mode

    def _spec_col(title, rows, x, W=_W):
        _section_hdr(slide, x, 1.68, W, 0.46, title)
        for i, (k, v) in enumerate(rows):
            yy = 2.22 + i * 0.54
            bg = _LGRY if i % 2 == 0 else _WHT
            _rect(slide, x, yy, W, 0.52, bg)
            _tb(slide, x+0.18, yy+0.10, W*0.50, 0.36, k, 9.5, color=_MGY)
            _tb(slide, x+W*0.52, yy+0.10, W*0.46, 0.36, v, 10,
                bold=True, color=_DGY, align="right")

    c_rate  = params.get("c_rate", 0.25)
    pv_rows = [
        ("Peak power",          _fmw(pv_kwp, "kWp")),
        ("System losses",       f"{params.get('pv_loss',14):.0f}%"),
        ("Tilt / Azimuth",      f"{params.get('tilt',20):.0f}° / "
                                f"{params.get('azimuth',180):.0f}°"),
        ("Annual generation",   _fmw(annual_kwh, "kWh/yr")),
        ("Specific yield",      f"{spec_yield:,.0f} kWh/kWp/yr"),
        ("Summer daily avg",    f"{pvgis_data.get('summer_daily_kwh',0):.1f} kWh/day"),
        ("Winter daily avg",    f"{pvgis_data.get('winter_daily_kwh',0):.1f} kWh/day"),
        ("Annual degradation",  f"{params.get('pv_degradation',0.5):.1f}% p.a."),
    ]
    bess_rows = [
        ("Usable capacity",      _fmw(bess_kh, "kWh")),
        ("Max charge/discharge", f"{_fmw(bess_kh*c_rate,'kW')}  "
                                 f"({params.get('c_rate_label','0.25C')})"),
        ("Depth of discharge",   f"{params.get('dod',90):.0f}%"),
        ("Round-trip efficiency",f"{params.get('rte',90):.0f}%"),
        ("Design cycle rate",    f"{params.get('bess_cycles',365):.0f} cycles/yr"),
        ("Tariff mode",          (params.get("tariff_mode","—")[:28]+"…")
                                 if len(params.get("tariff_mode",""))>28
                                 else params.get("tariff_mode","—")),
        ("Analysis period",      "20 years"),
        ("Dispatch strategy",    "Peak shaving + arbitrage"),
    ]

    if _dual:
        _spec_col("SOLAR PV ARRAY",         pv_rows,   0.28, _W)
        _spec_col("BATTERY ENERGY STORAGE", bess_rows, _x2,  _W)
    elif has_pv:
        _spec_col("SOLAR PV ARRAY",         pv_rows,   0.28, _W)
    else:
        _spec_col("BATTERY ENERGY STORAGE", bess_rows, 0.28, _W)

    # ── Right product image panel (single-component mode only) ────────────────
    # PV+BESS: spec tables already fill the full width — no panel added.
    # Single mode: 5.90"-wide product showcase to the right of the spec table.
    if not _dual and _am is not None:
        img_path, img_caption = _am.system_product_image(has_pv, has_bess)
        _PNL_X = 7.12; _PNL_Y = 1.68; _PNL_W = 5.90; _PNL_H = 5.48
        # White card background with subtle border
        _rect(slide, _PNL_X, _PNL_Y, _PNL_W, _PNL_H, _WHT, line=_SEP)
        # Thin Huawei-red accent at the top of the card
        _rect(slide, _PNL_X, _PNL_Y, _PNL_W, 0.05, _HRD)
        if img_path and os.path.exists(img_path):
            # Centre the product image with 0.30" padding on each side
            slide.shapes.add_picture(
                img_path,
                _in(_PNL_X + 0.35), _in(_PNL_Y + 0.22),
                _in(_PNL_W - 0.70), _in(_PNL_H - 0.66),
            )
        # Caption inside the card (bottom-aligned)
        _tb(slide, _PNL_X, _PNL_Y + _PNL_H - 0.40, _PNL_W, 0.30,
            img_caption, 8, color=_MGY, align="center")

    # ── Bottom separator + PVGIS note ─────────────────────────────────────────
    # In single mode limit the separator/narrative to the spec table width only
    _sep_w = 12.78 if _dual else 6.60
    _rect(slide, 0.28, 6.78, _sep_w, 0.03, _SEP)
    _narrative(slide, 0.28, 6.85, _sep_w, 0.28,
               "Irradiance sourced from EU PVGIS API — crystalline silicon, "
               "free-mounted, site coordinates.  Fallback: 1,650 kWh/kWp/yr.")
    _footer(slide, company, page, total=total)


def _s4_financial(prs, results: dict, fin_df, company: str,
                  page: int = 4, total: int = 11):
    """Slide 4 – Financial Returns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    npv     = results.get("npv", 0) or 0
    irr     = results.get("irr", 0) or 0
    payback = results.get("payback") or 0
    capex   = results.get("total_capex", 0) or 0

    def _m(v): return f"R {abs(v)/1e6:.2f}M"
    sign = "+" if npv >= 0 else "−"
    headline = (f"Project returns {sign}R {abs(npv)/1e6:.1f}M NPV at {irr:.1f}% IRR "
                f"— full payback by Year {payback:.0f}")
    _header_bar(slide, headline,
                f"CAPEX {_m(capex)}  ·  Discount rate "
                f"{results.get('discount_rate',12):.1f}%  ·  "
                f"CIT {results.get('tax_rate',27):.0f}%  ·  "
                "Section 12B: 100% Year-1 (<1 MW) or 50/30/20% (>1 MW) applied")

    chart_png = _png_cashflow(fin_df) if fin_df is not None else None
    if chart_png:
        slide.shapes.add_picture(io.BytesIO(chart_png),
                                 _in(0.28), _in(1.68), _in(8.80), _in(3.62))
    else:
        _rect(slide, 0.28, 1.68, 8.80, 3.62, _LGRY)
        _tb(slide, 1.0, 3.30, 7.0, 0.45,
            "Run simulation first to generate cash-flow chart",
            11, color=_MGY, align="center")

    _tb(slide, 0.28, 5.38, 8.80, 0.28,
        "Red bars = positive NCF  ·  Bright-red = negative  ·  "
        "Black line = cumulative  ·  Dashed = break-even year",
        8, color=_MGY, align="center")

    for i, (lbl, val, sub) in enumerate([
        ("CAPEX",       _m(capex),              "Total investment"),
        ("20-yr NPV",   f"{sign}{_m(abs(npv))}", "After tax, discounted"),
        ("IRR",         f"{irr:.1f}%",           "Unlevered, after-tax"),
        ("Payback",     f"Year {payback:.1f}",   "Simple CF basis"),
    ]):
        _kpi_block(slide, 9.55, 1.68+i*1.36, lbl, val, sub, w=3.20)

    _rect(slide, 0.28, 5.72, 12.78, 0.03, _SEP)
    _narrative(slide, 0.28, 5.82, 12.78, 0.62,
               "The 20-year model includes O&M escalation, PV degradation, and "
               "BESS SOH decline.  Section 12B (100% Year-1 for <1 MW assets) "
               "materially improves early cash flows.  NPV remains positive across "
               "±2% tariff-escalation sensitivity (available in the full Excel model).")
    _footer(slide, company, page, total=total)


def _s4b_financial_table(prs, results: dict, fin_df, company: str,
                          page: int = 5, total: int = 12):
    """Slide 4b — 20-Year Financial Projections Table (two-panel layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    capex   = results.get("total_capex", 0) or 0
    tax_r   = results.get("tax_rate", 27)
    disc_r  = results.get("discount_rate", 12)

    _header_bar(
        slide,
        "20-Year Projected Cash Flow — Year-by-Year Breakdown",
        (f"CAPEX R {abs(capex)/1e6:.2f}M  ·  All values in R millions  ·  "
         f"Discount {disc_r:.1f}%  ·  CIT {tax_r:.0f}%  ·  "
         "Green Cum.CF = post-payback positive territory"),
    )

    if fin_df is None or len(fin_df) == 0:
        _narrative(slide, 0.28, 3.50, 12.78, 0.60,
                   "Run simulation first to generate the 20-year financial table.",
                   )
        _footer(slide, company, page, total=total)
        return

    # ── Layout constants ──────────────────────────────────────────────────────
    TBL_Y   = 1.68          # top of table area
    TBL_H   = 5.54          # y=1.68 → y=7.22 (above footer)
    HDR_H   = 0.38          # header row height
    ROW_H   = (TBL_H - HDR_H) / 10  # 10 rows per panel
    GAP     = 0.18          # gap between panels

    # Column defs per panel: (header_text, width_in, data_key_or_fn, align)
    #   data_key_or_fn: either a column name in fin_df, or a callable(row)->str
    def _rv(v, neg_brackets=False):
        """Format R-million value."""
        if abs(v) >= 100:
            s = f"{v/1e6:.1f}"
        elif abs(v) >= 10:
            s = f"{v/1e6:.2f}"
        else:
            s = f"{v/1e6:.2f}"
        if neg_brackets and v < 0:
            return f"({abs(v)/1e6:.2f})"
        return s

    # payback year (rounded up to next integer year)
    payback = results.get("payback") or 0
    payback_yr = int(payback) + (1 if payback % 1 > 0.01 else 0)

    PANEL_COLS = [
        # header_text       width  key / lambda                           align
        ("Yr",              0.42,  lambda r: str(int(r["Year"])),         "center"),
        ("Gross Saving\nR M", 1.60, lambda r: _rv(r.get("Total Saving (ZAR)", 0)), "right"),
        ("NCF\nR M",        1.60,  lambda r: _rv(r.get("Net Cash Flow NCF (ZAR)", 0)), "right"),
        ("Cum. CF\nR M",    2.22,  lambda r: _rv(r.get("Cumulative CF (ZAR)", 0)), "right"),
        ("",                0.36,  lambda r: "✓" if r.get("Status", "✓") == "✓" else "EoL", "center"),
    ]
    L_PAN_W = sum(c[1] for c in PANEL_COLS)  # 6.20"
    R_PAN_W = L_PAN_W + 0.20                 # 6.40" (slightly wider)
    # recalc R panel col widths (proportional scale)
    _scale  = R_PAN_W / L_PAN_W
    R_COLS  = [(h, w * _scale, fn, a) for h, w, fn, a in PANEL_COLS]

    def _draw_panel(x_start, pan_w, col_defs, year_start, year_end):
        rows_slice = fin_df[
            (fin_df["Year"] >= year_start) & (fin_df["Year"] <= year_end)
        ]
        n = len(rows_slice)

        # compute column x positions
        xs = []
        cx = x_start
        for _, w, _, _ in col_defs:
            xs.append(cx)
            cx += w

        # ── Header row ────────────────────────────────────────────────────
        _rect(slide, x_start, TBL_Y, pan_w, HDR_H, _DGY)
        for (hdr, w, _, align), hx in zip(col_defs, xs):
            _tb(slide, hx + 0.03, TBL_Y + 0.05, w - 0.06, HDR_H - 0.06,
                hdr, 7.5, bold=True, color=_WHT, align=align, wrap=False)

        # ── Data rows ─────────────────────────────────────────────────────
        for ri, (_, row) in enumerate(rows_slice.iterrows()):
            yr     = int(row["Year"])
            ry     = TBL_Y + HDR_H + ri * ROW_H
            is_eol = (str(row.get("Status", "✓")) != "✓")
            cum_cf = row.get("Cumulative CF (ZAR)", 0)

            # Row background
            if yr == payback_yr:
                bg = _GLD
            elif ri % 2 == 0:
                bg = _LGRY
            else:
                bg = _WHT
            _rect(slide, x_start, ry, pan_w, ROW_H - 0.01, bg)

            for ci, ((_, w, fn, align), hx) in enumerate(zip(col_defs, xs)):
                try:
                    val_str = fn(row)
                except Exception:
                    val_str = "—"

                # Text colour
                if ci == 3:        # Cum. CF column
                    if cum_cf > 0:
                        clr = "15803D"   # dark green
                    else:
                        clr = _NEG
                elif ci == 4:      # Status column
                    clr = "15803D" if not is_eol else _HRD
                else:
                    clr = _MGY if is_eol else _DGY

                _tb(slide, hx + 0.04, ry + 0.04,
                    w - 0.08, ROW_H - 0.06,
                    val_str, 8.5, bold=(ci == 0),
                    color=clr, align=align, wrap=False)

    # Left panel: years 1-10
    _draw_panel(0.28, L_PAN_W, PANEL_COLS, 1, 10)
    # Right panel: years 11-20
    _draw_panel(0.28 + L_PAN_W + GAP, R_PAN_W, R_COLS, 11, 20)

    # Panel labels above the header rows
    _tb(slide, 0.28,              TBL_Y - 0.26, L_PAN_W, 0.24,
        "◀  Years 1 – 10", 8.5, bold=True, color=_HRD)
    _tb(slide, 0.28 + L_PAN_W + GAP, TBL_Y - 0.26, R_PAN_W, 0.24,
        "◀  Years 11 – 20", 8.5, bold=True, color=_HRD)

    _footer(slide, company, page, total=total)


def _s5_energy(prs, pvgis_data: dict, results: dict, company: str,
               page: int = 5, total: int = 11, has_bess: bool = True):
    """Slide 5 – Energy Analysis (PV yield + optional BESS dispatch)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    d1          = results.get("dispatch_yr1") or {}
    annual_gen  = pvgis_data.get("annual_kwh", 0)
    annual_load = d1.get("annual_load_kWh", 0)
    annual_grid = d1.get("annual_grid_buy_kWh", 0)
    annual_disc = d1.get("annual_discharge_kWh", 0)
    self_suf    = (1-annual_grid/annual_load)*100 if annual_load else 0
    displaced   = annual_load-annual_grid if annual_load else 0

    headline = (f"Year-1 PV generation {annual_gen/1e3:,.0f} MWh — "
                f"{self_suf:.0f}% self-sufficiency, "
                f"{displaced/1e3:,.0f} MWh of grid purchases displaced")
    _header_bar(slide, headline,
                "June–August winter peak (darker bars) coincides with highest "
                "Eskom TOU rates — battery dispatches aggressively in this window.")

    chart_png = _png_monthly(pvgis_data)
    if chart_png:
        slide.shapes.add_picture(io.BytesIO(chart_png),
                                 _in(0.28), _in(1.68), _in(8.80), _in(3.18))
    else:
        _rect(slide, 0.28, 1.68, 8.80, 3.18, _LGRY)

    kpi_rows = [
        ("PV GENERATION",  f"{annual_gen/1e3:,.0f} MWh",  "Year-1 total"),
        ("SITE LOAD",      f"{annual_load/1e3:,.0f} MWh" if annual_load else "—",
                           "Yr-1 total demand"),
        ("GRID PURCHASE",  f"{annual_grid/1e3:,.0f} MWh" if annual_grid else "—",
                           "After PV" + ("+BESS" if has_bess else "")),
    ]
    if has_bess:
        kpi_rows.append(("BESS DISCHARGE",
                         f"{annual_disc/1e3:,.0f} MWh" if annual_disc else "—",
                         "Battery output"))
    kpi_h = 5.04 / len(kpi_rows)
    for i, (lbl, val, sub) in enumerate(kpi_rows):
        _kpi_block(slide, 9.55, 1.68+i*kpi_h, lbl, val, sub, w=3.20)

    # Monthly mini-table
    _rect(slide, 0.28, 4.96, 8.80, 0.03, _SEP)
    months  = ["Jan","Feb","Mar","Apr","May","Jun",
               "Jul","Aug","Sep","Oct","Nov","Dec"]
    monthly = pvgis_data.get("monthly_kwh", [0]*12)
    cw = 8.74/12
    for j, (m, v) in enumerate(zip(months, monthly)):
        xx  = 0.31 + j*cw
        bg  = _LRD if j in (5,6,7) else _WHT
        clr = _DRD if j in (5,6,7) else _DGY
        _rect(slide, xx, 5.02, cw-0.04, 0.68, bg)
        _tb(slide, xx+0.02, 5.05, cw-0.06, 0.27, m, 7.5,
            color=_MGY, align="center")
        _tb(slide, xx+0.02, 5.34, cw-0.06, 0.30, f"{v/1000:.1f}k", 8.5,
            bold=True, color=clr, align="center")

    _rect(slide, 0.28, 5.82, 12.78, 0.03, _SEP)
    if has_bess:
        _narrative(slide, 0.28, 5.92, 12.78, 0.65,
                   f"Self-sufficiency of {self_suf:.0f}% means only "
                   f"{annual_grid/1e3:,.0f} MWh is drawn from the grid in Year 1.  "
                   "The BESS charges during off-peak periods and discharges into "
                   "morning (07:00–09:00) and evening (17:00–20:00) Eskom peak "
                   "windows — shaded red columns indicate the high-demand season.")
    else:
        _narrative(slide, 0.28, 5.92, 12.78, 0.65,
                   f"Solar PV generates {annual_gen/1e3:,.0f} MWh in Year 1, "
                   f"displacing {displaced/1e3:,.0f} MWh of grid purchases and "
                   f"achieving {self_suf:.0f}% self-sufficiency.  "
                   "Shaded red columns indicate the high-demand winter season "
                   "(June–August) where PV self-consumption delivers peak savings.")
    _footer(slide, company, page, total=total)


def _s6_tariff(prs, params: dict, results: dict, company: str,
               page: int = 6, total: int = 11):
    """Slide 6 – Tariff Opportunity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w_pk  = params.get("w_evening_peak", 8.13)
    w_op  = params.get("w_off_peak", 1.57)
    esc   = params.get("tariff_escalation", 8.0)
    ratio = w_pk/w_op if w_op else 0

    headline = (f"Winter peak R {w_pk:.2f}/kWh vs off-peak R {w_op:.2f}/kWh "
                f"({ratio:.1f}×) — R {w_pk-w_op:.2f}/kWh BESS arbitrage spread")
    _header_bar(slide, headline,
                f"Tariff: {params.get('tariff_mode','—')}  ·  "
                f"Escalation {esc:.1f}% p.a.  ·  All rates incl. 15% VAT")

    # TOU table
    _section_hdr(slide, 0.28, 1.68, 5.80, 0.44, "TOU ENERGY RATES  (ZAR / kWh)")
    rate_rows = [
        ("Morning Peak",
         f"R {params.get('w_morning_peak',0):.4f}",
         f"R {params.get('s_morning_peak',0):.4f}"),
        ("Evening Peak",
         f"R {params.get('w_evening_peak',0):.4f}",
         f"R {params.get('s_evening_peak',0):.4f}"),
        ("Standard",
         f"R {params.get('w_standard',0):.4f}",
         f"R {params.get('s_standard',0):.4f}"),
        ("Off-Peak",
         f"R {params.get('w_off_peak',0):.4f}",
         f"R {params.get('s_off_peak',0):.4f}"),
    ]
    # Column headers
    for cx, hdr in zip([0.44,2.28,4.14], ["Period","Winter","Summer"]):
        _tb(slide, cx, 2.18, 1.8, 0.28, hdr, 8.5, bold=True, color=_MGY)
    for i, (a, b, c) in enumerate(rate_rows):
        yy = 2.50+i*0.58
        bg = _LGRY if i % 2 == 0 else _WHT
        _rect(slide, 0.28, yy, 5.80, 0.56, bg)
        _tb(slide, 0.44, yy+0.10, 1.7, 0.36, a, 9.5, color=_MGY)
        _tb(slide, 2.28, yy+0.10, 1.6, 0.36, b, 10,
            bold=True, color=_DGY, align="center")
        _tb(slide, 4.14, yy+0.10, 1.8, 0.36, c, 10,
            bold=True, color=_DGY, align="center")

    # ── Dispatch logic block — split: narrative (left 55%) + BESS product (right 45%) ──
    _rect(slide, 0.28, 4.96, 5.80, 1.58, _LRD)
    _rect(slide, 0.28, 4.96, 0.07, 1.58, _HRD)
    _txt_w = 3.10   # text area width (narrowed to make room for product image)
    _tb(slide, 0.44, 5.04, _txt_w, 0.32,
        "Dispatch Strategy", 10, bold=True, color=_DGY)
    _narrative(slide, 0.44, 5.40, _txt_w, 1.05,
               f"BESS charges from PV surplus and cheap off-peak grid power "
               f"(≤R {w_op:.2f}/kWh). It discharges during 07:00–09:00 and "
               f"17:00–20:00 peak windows — capturing the full "
               f"R {w_pk-w_op:.2f}/kWh arbitrage spread.")
    # BESS product image inset (right 40% of dispatch block)
    if _am is not None:
        _bess_img = _am.tariff_product_image()
        if _bess_img and os.path.exists(_bess_img):
            _rect(slide, 3.58, 5.01, 2.46, 1.50, _WHT)
            slide.shapes.add_picture(
                _bess_img, _in(3.62), _in(5.04), _in(2.38), _in(1.44))

    # Right: savings KPIs + escalation
    d1      = results.get("dispatch_yr1") or {}
    ann_sav = d1.get("annual_savings_zar", 0) or 0
    def _mv(v): return (f"R {abs(v)/1e6:.2f}M/yr" if abs(v)>=1e6
                        else f"R {abs(v)/1e3:.0f}k/yr")

    _kpi_block(slide, 6.55, 1.68, "YR-1 GRID SAVINGS",
               _mv(ann_sav) if ann_sav else "Run sim",
               "Avoided grid purchase cost", w=3.20)
    _kpi_block(slide, 10.0, 1.68, "TARIFF ESCALATION",
               f"{esc:.1f}% p.a.", "Annual increase assumed", w=3.05)

    # Escalation projection
    _section_hdr(slide, 6.55, 3.22, 6.50, 0.40,
                 "Peak rate projection (illustrative)")
    base = w_pk
    for i, yr in enumerate([1, 3, 5, 10, 15, 20]):
        proj = base * (1+esc/100)**(yr-1)
        xx   = 6.58+i*1.08
        bg   = _LRD if yr >= 10 else _LGRY
        clr  = _HRD if yr >= 10 else _DGY
        _rect(slide, xx, 3.66, 1.05, 1.02, bg)
        _tb(slide, xx+0.05, 3.72, 0.95, 0.28,
            f"Yr {yr}", 7.5, color=_MGY, align="center")
        _tb(slide, xx+0.05, 4.02, 0.95, 0.50,
            f"R{proj:.2f}", 10.5, bold=True, color=clr, align="center")

    _narrative(slide, 6.55, 5.72, 6.50, 0.62,
               f"At {esc:.1f}% p.a. escalation, the evening peak rate reaches "
               f"≈R {base*(1+esc/100)**9:.2f}/kWh by Year 10 and "
               f"≈R {base*(1+esc/100)**19:.2f}/kWh by Year 20.  "
               "Each R 0.10/kWh of additional escalation adds meaningfully to "
               "project NPV — early commissioning locks in these economics for "
               "the full 20-year analysis period.")
    _footer(slide, company, page, total=total)


def _s7_roadmap(prs, results: dict, params: dict, company: str,
                page: int = 7, total: int = 11):
    """Slide 7 – Implementation Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bess_lead = results.get("bess_lead_months",
                            params.get("bess_lead_months", 6))
    pv_lead   = results.get("pv_lead_months",
                            params.get("pv_lead_months", 12))

    headline = (f"BESS commissioned Month {bess_lead}, "
                f"full PV+BESS live Month {pv_lead} — "
                "phased delivery de-risks procurement and accelerates returns")
    _header_bar(slide, headline,
                "Timelines indicative — subject to site survey, grid connection "
                "application, and equipment availability.")

    total_months = pv_lead + 2
    bar_y = 2.00; bar_h = 0.65; bar_w = 12.50

    phases = [
        (0,         3,          "Feasibility &\nPermitting",   _MGY),
        (3,         bess_lead,  "Procurement\n& Civil Works",  _DGY),
        (bess_lead, pv_lead,    "BESS\nCommissioned",          _HRD),
        (pv_lead,   pv_lead+3,  "Full System\nOperational",    _AMBR),
    ]
    _rect(slide, 0.28, bar_y, bar_w, bar_h, _LGRY)
    for m_s, m_e, label, clr in phases:
        if m_e <= m_s: continue
        x = 0.28 + (m_s/total_months)*bar_w
        w = ((m_e-m_s)/total_months)*bar_w
        _rect(slide, x, bar_y, w-0.03, bar_h, clr)
        if w > 0.9:
            _tb(slide, x+0.08, bar_y+0.06, w-0.12, bar_h-0.10,
                label, 7.5, bold=True, color=_WHT)

    for m in range(0, total_months+1, 3):
        x = 0.28 + (m/total_months)*bar_w
        _rect(slide, x, bar_y+bar_h, 0.01, 0.14, _SEP)
        _tb(slide, x-0.20, bar_y+bar_h+0.16, 0.5, 0.24,
            f"M{m}", 7.5, color=_MGY, align="center")

    for m_num, label in [
        (bess_lead, f"BESS Live  M{bess_lead}"),
        (pv_lead,   f"PV Live  M{pv_lead}"),
    ]:
        x = 0.28 + (m_num/total_months)*bar_w
        _tb(slide, x-0.25, bar_y-0.52, 2.2, 0.38,
            label, 9, bold=True, color=_HRD)
        _rect(slide, x, bar_y-0.14, 0.03, 0.16, _HRD)

    _section_hdr(slide, 0.28, 3.38, 12.50, 0.40, "RECOMMENDED NEXT STEPS")

    steps = [
        ("1. Site Assessment",
         "Commission shading analysis, roof/ground survey, and single-line "
         "diagram to confirm design assumptions."),
        ("2. Grid Application",
         "Submit a Generation Connection Application (GCA) to the network "
         "operator — typical lead time 4–8 weeks."),
        ("3. Contractor RFQ",
         "Issue RFQ to ≥3 EPC contractors with South African PV+BESS "
         "references.  Specify Huawei FusionSolar / LUNA2000 or equivalent."),
        ("4. Financial Close",
         "Finalise funding structure (equity, debt, or PPA/RESCO). Section 12B "
         "requires asset ownership before 31 March tax year-end."),
    ]
    for i, (hdr, body) in enumerate(steps):
        col, row = i % 2, i // 2
        x = 0.28+col*6.40; y = 3.88+row*1.55
        _rect(slide, x, y, 6.20, 1.46, _LRD if row == 0 else _LGRY)
        _rect(slide, x, y, 0.07, 1.46, _HRD)
        _tb(slide, x+0.18, y+0.10, 5.85, 0.30,
            hdr, 10, bold=True, color=_DGY)
        _narrative(slide, x+0.18, y+0.44, 5.85, 0.95, body)

    _footer(slide, company, page, total=total)


def _s8_market(prs, company: str, page: int = 10, total: int = 11):
    """Slide – SA Eskom Megaflex Tariff Growth Trend (10-year history + projection)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _header_bar(
        slide,
        "SA Eskom Megaflex tariff has compounded at +10.8% p.a. over 10 years — BTM economics strengthen every year",
        "Source: NERSA tariff determinations FY2016–FY2027  ·  Representative blended effective rate, VAT-inclusive  "
        "·  Projections at NERSA FY2027-approved +8.83% forward",
    )

    # ── 4 KPI cards ──────────────────────────────────────────────────────────
    kpis = [
        ("10-YR CAGR",       "+10.8% p.a.",  "FY2016 → FY2026\nEskom Megaflex"),
        ("FY2026 RATE",      "≈R 2.56/kWh",  "NERSA-approved +8.76%\nvs FY2025"),
        ("EST. FY2030 RATE", "≈R 3.55/kWh",  "At 8.5% p.a. forward\n+R 0.99/kWh in 4 yrs"),
        ("CUMUL. INCREASE",  "+184%",         "FY2016 → FY2026\n(R 0.90 → R 2.56/kWh)"),
    ]
    for i, (lbl, val, sub) in enumerate(kpis):
        _kpi_block(slide, 0.28 + i * 3.21, 1.68, lbl, val, sub, w=3.08)

    # ── Separator ─────────────────────────────────────────────────────────────
    _rect(slide, 0.28, 3.12, 12.78, 0.03, _SEP)

    # ── Tariff trend chart ────────────────────────────────────────────────────
    chart_png = _png_megaflex_tariff()
    if chart_png:
        slide.shapes.add_picture(io.BytesIO(chart_png),
                                 _in(0.28), _in(3.22), _in(12.78), _in(3.72))
    else:
        _rect(slide, 0.28, 3.22, 12.78, 3.72, _LGRY)
        _tb(slide, 0.28, 4.90, 12.78, 0.40,
            "Megaflex tariff chart unavailable — run pip install matplotlib",
            10, color=_MGY, align="center")

    # ── Bottom source note ────────────────────────────────────────────────────
    _tb(slide, 0.28, 7.04, 12.78, 0.16,
        "NERSA approved an 8.76% increase for April 2026 and 8.83% for April 2027.  "
        "At compound escalation, a site paying ≈R 2.56/kWh today will pay ≈R 3.55/kWh by 2030 "
        "— an additional 39% cost increase in just 4 years.",
        7.5, color=_MGY, wrap=True)

    _footer(slide, company, page, total=total)


def _s9_assumptions(prs, params: dict, company: str,
                    page: int = 11, total: int = 11):
    """Slide 9 – Assumptions & Disclaimer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, 13.33, 7.5, _BLK)         # black background
    _rect(slide, 0, 0, 0.12, 7.5, _HRD)           # Huawei red left strip
    _rect(slide, 0.28, 0.55, 12.78, 5.72, _WHT)   # white content panel

    _tb(slide, 0.28, 0.04, 10.0, 0.48,
        "Key Assumptions  &  Disclaimer", 19, bold=True, color=_WHT)

    assumptions = [
        ("PV degradation",        f"{params.get('pv_degradation',0.5):.1f}% per year (linear)"),
        ("Round-trip efficiency",  f"{params.get('rte',90):.0f}%  (AC-AC)"),
        ("Analysis horizon",       "20 years from PV commissioning"),
        ("Tariff escalation",      f"{params.get('tariff_escalation',8.0):.1f}% per annum (real)"),
        ("Discount rate (WACC)",   f"{params.get('discount_rate',12.0):.1f}%  (nominal, after-tax)"),
        ("Corporate Income Tax",   f"{params.get('tax_rate',27):.0f}%  (CIT)"),
        ("Section 12B",            "100% Year-1 for <1 MW; 50/30/20% over 3 yrs for >1 MW"),
        ("Battery DoD",            f"{params.get('dod',90):.0f}%  (usable capacity)"),
        ("FX rate",                f"USD/ZAR  {params.get('forex_usd_zar',18.5):.2f}"),
        ("Irradiance data",        "EU PVGIS API — crystSi, free-mount, site lat/lon"),
    ]
    mid = len(assumptions)//2
    for col_i, chunk in enumerate([assumptions[:mid], assumptions[mid:]]):
        x = 0.42+col_i*6.45
        _section_hdr(slide, x-0.08, 0.58, 6.28, 0.38,
                     "TECHNICAL / FINANCIAL" if col_i == 0
                     else "REGULATORY / MARKET")
        for i, (k, v) in enumerate(chunk):
            yy = 1.04+i*0.54
            bg = _LGRY if i % 2 == 0 else _WHT
            _rect(slide, x-0.08, yy, 6.28, 0.52, bg)
            _tb(slide, x+0.08, yy+0.10, 2.80, 0.36, k, 9, color=_MGY)
            _tb(slide, x+2.92, yy+0.10, 3.20, 0.36, v, 9.5,
                bold=True, color=_DGY)

    # Disclaimer
    _rect(slide, 0.28, 6.32, 12.78, 0.42, _HRD)
    _tb(slide, 0.44, 6.37, 12.4, 0.30,
        "DISCLAIMER — READ BEFORE ACTING ON THIS DOCUMENT",
        8.5, bold=True, color=_WHT)
    disc = (
        "Prepared solely for informational and decision-support purposes. "
        "Outputs are based on modelled assumptions; actual performance, tariff "
        "levels, and regulatory conditions may differ materially.  Not investment "
        "advice.  All tariff figures inclusive of 15% South African VAT.  "
        "Obtain independent engineering, legal, and financial advice before "
        "committing capital.  "
        f"© {date.today().year} {company}.  "
        f"Solution Info: {_SOL_URL}"
    )
    _tb(slide, 0.28, 6.80, 12.78, 0.60, disc, 8, color="BBBBBB", wrap=True)
    _footer(slide, company, page, total=total)


def _s8_huawei_partner(prs, company: str, page: int = 8, total: int = 11):
    """Slide 8 — Huawei Digital Power · SA Partner Credentials."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ASSETS = os.path.join(os.path.dirname(__file__), "assets")

    _header_bar(
        slide,
        "Huawei Digital Power — Proven Technology Partner in South Africa",
        "400+ C&I PV+BESS installations · 1,400+ SA employees · "
        "BEE Level 2 supplier · Founded Johannesburg 1998",
    )

    # ── Left dark panel (38%) — SA campus photo + company facts ──────────────
    _rect(slide, 0.28, 0.68, 4.90, 6.54, _DNAV)          # dark navy panel bg

    campus_img = os.path.join(ASSETS, "hw_sa_grid.jpg")
    if os.path.exists(campus_img):
        slide.shapes.add_picture(
            campus_img, _in(0.28), _in(0.68), _in(4.90), _in(2.78))
    else:
        _rect(slide, 0.28, 0.68, 4.90, 2.78, "1C2A3A")
        _tb(slide, 0.28, 1.88, 4.90, 0.40,
            "HUAWEI South Africa Campus · Johannesburg",
            9, color="7A8FA8", align="center")

    # Red rule under photo
    _rect(slide, 0.28, 3.46, 4.90, 0.04, _HRD)

    # SA facts heading
    _tb(slide, 0.44, 3.56, 4.60, 0.28,
        "HUAWEI SOUTH AFRICA", 8, bold=True, color=_GLD)

    sa_facts = [
        ("Founded",        "1998 · Johannesburg HQ"),
        ("Revenue",        "US$1 Billion per year"),
        ("Staff",          "1,400+ employees · BEE Level 2"),
        ("Population",     "Serves 80% of SA population"),
        ("Utility PV",     "20+ plants incl. Scatec 540 MW RMIPP"),
        ("C&I / Res",      "400+ C&I sites · 10,000+ residential"),
    ]
    for k, (label, val) in enumerate(sa_facts):
        yy = 3.94 + k * 0.50
        _tb(slide, 0.44,  yy, 1.70, 0.38, label.upper(),
            7,   bold=True, color="7A8FA8")
        _tb(slide, 2.18, yy, 2.85, 0.38, val,
            8.5, bold=False, color="D1D5DB")

    # ── Right panel — global stats (3 rows × 3 KPIs each) ────────────────────
    col_x = [5.42, 7.86, 10.30]
    card_w = 2.32

    # Row 1: GLOBAL FOOTPRINT
    _section_hdr(slide, 5.42, 0.68, 7.63, 0.38, "GLOBAL FOOTPRINT")
    for k, (lbl, val, sub) in enumerate([
        ("EMPLOYEES",      "213,000",    "worldwide"),
        ("COUNTRIES",      "170+",       "countries & regions"),
        ("R&D INVESTMENT", "15%+",       "of annual revenue"),
    ]):
        x = col_x[k]
        _rect(slide, x, 1.14, card_w, 1.40, _DNAV)
        _tb(slide, x+0.14, 1.21, card_w-0.20, 0.26, lbl, 7, bold=True, color="9CA3AF")
        _tb(slide, x+0.14, 1.47, card_w-0.20, 0.60, val, 26, bold=True, color=_GLD)
        _tb(slide, x+0.14, 2.09, card_w-0.20, 0.24, sub, 7, color="6B7280")

    # Row 2: R&D & TECHNOLOGY
    _section_hdr(slide, 5.42, 2.66, 7.63, 0.38, "R&D & TECHNOLOGY")
    for k, (lbl, val, sub) in enumerate([
        ("R&D CENTRES",    "13",         "global locations"),
        ("ACTIVE PATENTS", "3,484",      "in energy domain"),
        ("INSTALLED BASE", "800+ GW",    "PV inverter capacity"),
    ]):
        x = col_x[k]
        _rect(slide, x, 3.12, card_w, 1.40, _DNAV)
        _tb(slide, x+0.14, 3.19, card_w-0.20, 0.26, lbl, 7, bold=True, color="9CA3AF")
        _tb(slide, x+0.14, 3.45, card_w-0.20, 0.60, val, 26, bold=True, color=_GLD)
        _tb(slide, x+0.14, 4.07, card_w-0.20, 0.24, sub, 7, color="6B7280")

    # Row 3: GLOBAL IMPACT
    _section_hdr(slide, 5.42, 4.64, 7.63, 0.38, "GLOBAL IMPACT  (Dec 2025)")
    for k, (lbl, val, sub) in enumerate([
        ("GREEN POWER",    "2,087B kWh", "generated for customers"),
        ("CO₂ REDUCED", "1.06B tons", "carbon emissions avoided"),
        ("TREES EQUIV.",   "1.45B",      "billion trees equivalent"),
    ]):
        x = col_x[k]
        _rect(slide, x, 5.10, card_w, 2.10, _DNAV)
        _tb(slide, x+0.14, 5.17, card_w-0.20, 0.26, lbl, 7, bold=True, color="9CA3AF")
        _tb(slide, x+0.14, 5.43, card_w-0.20, 0.72, val, 22, bold=True, color=_GLD)
        _tb(slide, x+0.14, 6.17, card_w-0.20, 0.60, sub, 7, color="6B7280", wrap=True)

    _footer(slide, company, page, total=total)


def _s9_sa_projects(prs, company: str, page: int = 9, total: int = 11):
    """Slide 9 — South Africa Reference Projects  (2 × 2 photo-card grid)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ASSETS = os.path.join(os.path.dirname(__file__), "assets")

    _header_bar(
        slide,
        "Delivered at Scale — South Africa Reference Projects",
        "Proven BTM PV+BESS performance across C&I, commercial and agricultural sectors",
    )

    # ── 4 project cards: (image, title, location, specs_gold, bullets_white) ──
    projects = [
        ("hw_office_park.jpg",
         "Huawei Office Park",
         "Johannesburg, South Africa",
         "3,200 kW  /  6,400 kWh BESS",
         "32 × LUNA2000 200kWh units\n"
         "Reduces 7 tonnes diesel per day\n"
         "Seamless switching · COD Nov 2022"),

        ("hw_fresh_market.jpg",
         "Cape Town Fresh Market",
         "Cape Town, South Africa",
         "1.2 MW PV  +  3 → 12 MWh ESS",
         "Uninterrupted cold-chain power\n"
         "Multi-phase ESS expansion\n"
         "Continuous business operation"),

        ("hw_biru.jpg",
         "Biru Restaurant",
         "Johannesburg, South Africa",
         "200 kW  /  400 kWh BESS",
         "ROI = 5.34 years\n"
         "700 kWh/day green generation\n"
         "50 t CO₂ reduction · SME proof-point"),

        ("hw_winery.jpg",
         "Waterkloof Wine Estate",
         "Somerset West, Western Cape",
         "200 kW Rooftop PV",
         "1,100 kWh/day green generation\n"
         "'From Nature to Nature' commitment\n"
         "Agricultural & rural deployment"),
    ]

    # Card grid: 2 columns × 2 rows
    cell_w = 6.33
    cell_h = 3.05
    x_pos  = [0.28, 6.72]
    y_pos  = [0.68, 3.83]

    for idx, (img_file, title, location, specs, bullets) in enumerate(projects):
        col = idx % 2
        row = idx // 2
        cx, cy = x_pos[col], y_pos[row]

        # Dark base (visible when no photo or as fallback)
        _rect(slide, cx, cy, cell_w, cell_h, _DNAV)

        # Project photo (full-card)
        img_path = os.path.join(ASSETS, img_file)
        if os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path, _in(cx), _in(cy), _in(cell_w), _in(cell_h))

        # Solid dark strip at bottom (text area, ~42% of card height)
        strip_h = 1.28
        strip_y = cy + cell_h - strip_h
        shp = slide.shapes.add_shape(
            1, _in(cx), _in(strip_y), _in(cell_w), _in(strip_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb("080E18")
        shp.line.fill.background()

        # Red top accent strip on card
        _rect(slide, cx, cy, cell_w, 0.05, _HRD)

        # Location (small, muted)
        _tb(slide, cx+0.15, strip_y+0.06, cell_w-0.22, 0.24,
            location.upper(), 6.5, bold=False, color="7A8FA8")

        # Title (white bold)
        _tb(slide, cx+0.15, strip_y+0.28, cell_w-0.22, 0.38,
            title, 13, bold=True, color=_WHT)

        # Gold specs
        _tb(slide, cx+0.15, strip_y+0.66, cell_w-0.22, 0.30,
            specs, 10, bold=True, color=_GLD)

        # Bullet text (small white — only visible with enough space)
        # Tucked into top-left of card above the dark strip to avoid overflow
        bullet_y = cy + 0.18
        bullet_h = cell_h - strip_h - 0.24
        if bullet_h > 0.5:
            for bi, line in enumerate(bullets.split("\n")[:3]):
                _tb(slide, cx+0.15, bullet_y + bi*0.28, cell_w-0.22, 0.26,
                    f"· {line}", 7.5, color="D1D5DB")

    # Thin divider between rows
    _rect(slide, 0.28, 3.77, 12.78, 0.03, _SEP)

    _footer(slide, company, page, total=total)


# ── Public API ────────────────────────────────────────────────────────────────

def generate_pptx(
    params:          dict,
    results:         dict | None,
    fin_df,
    pvgis_data:      dict,
    project_name:    str = "",
    client_name:     str = "",
    consultant_name: str = "",
) -> bytes:
    """
    Generate an executive PPTX report (Huawei Digital Power style).
    Slide count auto-adapts to the configured PV / BESS scenario:
      • PV + BESS  → 11 slides (full deck)
      • PV only    → 10 slides (Tariff/BESS slide omitted)
      • BESS only  → 10 slides (Energy/PV yield slide omitted)

    Parameters
    ----------
    params          : session-state param dict (tariff rates, capex, etc.)
    results         : results dict (npv, irr, dispatch_yr1, lcoe, …)
    fin_df          : 20-year financial DataFrame (or None)
    pvgis_data      : PVGIS dict (annual_kwh, monthly_kwh, …)
    project_name    : cover title
    client_name     : client name — appears in "Confidential · Only For: …"
    consultant_name : EPC company on all slides (default "Energy Consulting")

    Returns bytes — pass directly to st.download_button(data=...)
    """
    from pptx import Presentation

    results    = results    or {}
    pvgis_data = pvgis_data or {}
    if not consultant_name:
        consultant_name = "Energy Consulting"
    company = consultant_name

    prs = Presentation()
    prs.slide_width  = _in(13.33)
    prs.slide_height = _in(7.5)

    pv_kwp   = params.get("pv_kwp",   0) or 0
    bess_kwh = params.get("bess_kwh", 0) or 0
    has_pv   = pv_kwp   > 0
    has_bess = bess_kwh > 0
    # Fallback: if both are zero (no simulation run) treat as PV+BESS
    if not has_pv and not has_bess:
        has_pv = has_bess = True

    # ── Dynamic slide list ────────────────────────────────────────────────────
    #   "financial_table" — 20-yr detail table, always included after "financial"
    #   "energy"  — PV yield analysis  → include only if has_pv
    #   "tariff"  — BESS arbitrage     → include only if has_bess
    _slide_ids = (["cover", "thesis", "system", "financial", "financial_table"]
                  + (["energy"] if has_pv   else [])
                  + (["tariff"] if has_bess else [])
                  + ["roadmap", "huawei_partner", "sa_projects", "market",
                     "assumptions"])
    _total = len(_slide_ids)

    for _pg, _sid in enumerate(_slide_ids, start=1):
        if _sid == "cover":
            _s1_cover(prs, project_name, client_name, consultant_name,
                      pv_kwp, bess_kwh, params.get("tariff_mode", ""),
                      has_pv=has_pv, has_bess=has_bess)
        elif _sid == "thesis":
            _s2_thesis(prs, results, params, company,
                       page=_pg, total=_total)
        elif _sid == "system":
            _s3_system(prs, params, pvgis_data, company,
                       page=_pg, total=_total,
                       has_pv=has_pv, has_bess=has_bess)
        elif _sid == "financial":
            _s4_financial(prs, results, fin_df, company,
                          page=_pg, total=_total)
        elif _sid == "financial_table":
            _s4b_financial_table(prs, results, fin_df, company,
                                 page=_pg, total=_total)
        elif _sid == "energy":
            _s5_energy(prs, pvgis_data, results, company,
                       page=_pg, total=_total, has_bess=has_bess)
        elif _sid == "tariff":
            _s6_tariff(prs, params, results, company,
                       page=_pg, total=_total)
        elif _sid == "roadmap":
            _s7_roadmap(prs, results, params, company,
                        page=_pg, total=_total)
        elif _sid == "huawei_partner":
            _s8_huawei_partner(prs, company, page=_pg, total=_total)
        elif _sid == "sa_projects":
            _s9_sa_projects(prs, company, page=_pg, total=_total)
        elif _sid == "market":
            _s8_market(prs, company, page=_pg, total=_total)
        elif _sid == "assumptions":
            _s9_assumptions(prs, params, company, page=_pg, total=_total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
